from .base import BaseAgent
from typing import Dict, Any
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from pptx import Presentation
from pathlib import Path

def _wrap(text: str, width: int = 90):
  # naive word-wrap for PDF lines
  words, line, n = text.split(), [], 0
  for w in words:
    if n + len(w) + (1 if line else 0) > width:
      yield " ".join(line); line, n = [w], len(w)
    else:
      line.append(w); n += len(w) + (1 if line[:-1] else 0)
  if line: yield " ".join(line)

class GenerationAgent(BaseAgent):
  name = "generator"

  def run(self, session_summary: str, out_dir: str, kind: str = "pdf") -> Dict[str, Any]:
    out = Path(out_dir); out.mkdir(parents=True, exist_ok=True)
    if kind == "pdf":
      path = out / "summary.pdf"
      c = canvas.Canvas(str(path), pagesize=A4)
      width, height = A4
      y = height - 72
      c.setFont("Helvetica-Bold", 16); c.drawString(72, y, "Video Analysis Summary"); y -= 24
      c.setFont("Helvetica", 11)
      for line in _wrap(session_summary, 90):
        if y < 72: c.showPage(); y = height - 72; c.setFont("Helvetica", 11)
        c.drawString(72, y, line); y -= 16
      c.save()
      return {"report": str(path)}
    else:
      path = out / "summary.pptx"
      prs = Presentation()
      s = prs.slides.add_slide(prs.slide_layouts[0])
      s.shapes.title.text = "Video Analysis Summary"
      s.placeholders[1].text = "Auto-generated (local)"
      s = prs.slides.add_slide(prs.slide_layouts[1])
      s.shapes.title.text = "Key Points"
      s.placeholders[1].text = session_summary
      prs.save(str(path))
      return {"report": str(path)}
